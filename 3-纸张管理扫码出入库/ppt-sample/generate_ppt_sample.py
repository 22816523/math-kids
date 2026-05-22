from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).with_name("AI参与原有功能迭代升级的产品设计经验_3页样稿.pptx")


NAVY = RGBColor(11, 30, 67)
DEEP_BLUE = RGBColor(17, 42, 93)
BLUE = RGBColor(46, 103, 248)
LIGHT_BLUE = RGBColor(220, 232, 255)
SOFT_BLUE = RGBColor(239, 245, 255)
PAPER = RGBColor(255, 253, 248)
SAND = RGBColor(247, 243, 234)
INK = RGBColor(19, 41, 79)
GRAY = RGBColor(97, 111, 137)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(255, 143, 61)
RED = RGBColor(239, 79, 69)
SOFT_RED = RGBColor(255, 244, 239)


def add_box(slide, left, top, width, height, text="", *,
            fill=None, line=None, radius=False,
            font_name="Microsoft YaHei", font_size=18,
            bold=False, color=INK, align=PP_ALIGN.LEFT,
            valign=MSO_ANCHOR.TOP, margin=0.12):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_text(slide, left, top, width, height, text, *,
             font_name="Microsoft YaHei", font_size=18,
             bold=False, color=INK, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, margin=0.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *,
                font_size=16, color=INK,
                spacing=0.12):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(7)
        p.bullet = True
    return box


def add_tag(slide, left, top, width, text, fill, font_color):
    return add_box(
        slide, left, top, width, Inches(0.34), text,
        fill=fill, line=None, radius=True,
        font_size=10, bold=True, color=font_color,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.03
    )


def add_slide_number(slide, current, total):
    add_box(
        slide, Inches(0.35), Inches(6.88), Inches(1.1), Inches(0.28),
        f"{current:02d} / {total:02d}",
        fill=PAPER, line=LIGHT_BLUE, radius=True,
        font_size=10, bold=True, color=INK,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02
    )


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, fill=RGBColor(238, 243, 251))
    add_box(slide, Inches(0.55), Inches(0.52), Inches(7.35), Inches(5.95), fill=PAPER, line=LIGHT_BLUE, radius=True)
    add_box(slide, Inches(8.1), Inches(0.52), Inches(4.68), Inches(3.78), fill=NAVY, line=None, radius=True)
    add_box(slide, Inches(8.1), Inches(4.48), Inches(4.68), Inches(2.0), fill=SOFT_BLUE, line=LIGHT_BLUE, radius=True)

    add_box(slide, Inches(0.9), Inches(0.9), Inches(1.45), Inches(0.36), "AI + PRODUCT", fill=PAPER, line=LIGHT_BLUE, radius=True,
            font_size=10, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)

    add_text(slide, Inches(0.95), Inches(1.55), Inches(6.4), Inches(1.7),
             "AI参与原有功能迭代升级的\n产品设计经验",
             font_name="Microsoft YaHei", font_size=26, bold=True, color=INK)
    add_text(slide, Inches(0.95), Inches(3.15), Inches(6.2), Inches(0.9),
             "以纸张管理扫码出入库项目为例，讲清楚 AI 在存量 ERP 功能迭代里如何给产品经理提效，以及如何避免失控。",
             font_size=15, color=GRAY)

    tags = [
        ("原有ERP迭代", Inches(0.95), Inches(4.45), Inches(1.2)),
        ("产品方法复盘", Inches(2.28), Inches(4.45), Inches(1.25)),
        ("AI协作提效", Inches(3.68), Inches(4.45), Inches(1.15)),
    ]
    for text, left, top, width in tags:
        add_tag(slide, left, top, width, text, WHITE, INK)

    add_text(slide, Inches(8.45), Inches(1.05), Inches(3.9), Inches(0.6),
             "核心判断", font_size=11, bold=True, color=RGBColor(191, 205, 236))
    add_text(slide, Inches(8.45), Inches(1.72), Inches(3.55), Inches(1.65),
             "AI 的价值，不在于替产品经理做判断，而在于更快地整理、对照、同步和暴露问题。",
             font_size=18, bold=True, color=WHITE)
    add_text(slide, Inches(8.45), Inches(3.45), Inches(3.65), Inches(0.35),
             "纸张管理扫码出入库项目 / 内部经验分享", font_size=10, color=RGBColor(205, 215, 239))

    stat_lefts = [Inches(8.35), Inches(9.93), Inches(11.51)]
    stats = [("1", "个存量 ERP 迭代案例"), ("7", "组可复用章节结构"), ("3+3", "踩坑与有效做法")]
    for left, (num, label) in zip(stat_lefts, stats):
      add_box(slide, left, Inches(4.86), Inches(1.25), Inches(1.18), fill=WHITE, line=LIGHT_BLUE, radius=True)
      add_text(slide, left + Inches(0.1), Inches(5.02), Inches(1.05), Inches(0.35), num, font_size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
      add_text(slide, left + Inches(0.04), Inches(5.45), Inches(1.16), Inches(0.42), label, font_size=8, color=GRAY, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 1, 3)


def slide_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, fill=RGBColor(238, 243, 251))
    add_box(slide, Inches(0.58), Inches(0.42), Inches(2.05), Inches(0.34), "CHAPTER 3", fill=PAPER, line=LIGHT_BLUE, radius=True,
            font_size=10, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)
    add_text(slide, Inches(0.65), Inches(0.98), Inches(5.2), Inches(0.8),
             "AI是怎么介入的，价值体现在哪", font_size=24, bold=True, color=INK)
    add_box(slide, Inches(10.95), Inches(0.42), Inches(1.85), Inches(0.34), "企业复盘型主模板", fill=PAPER, line=LIGHT_BLUE, radius=True,
            font_size=9, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)

    add_box(slide, Inches(0.65), Inches(1.7), Inches(4.35), Inches(4.95), fill=DEEP_BLUE, line=None, radius=True)
    add_text(slide, Inches(0.95), Inches(2.0), Inches(3.8), Inches(0.45), "不是少思考了，而是这些环节明显更快",
             font_size=18, bold=True, color=WHITE)

    value_items = [
        ("需求整理提效", "把跨多轮对话的零散信息，先收成结构草稿。"),
        ("差异核查提效", "把原始需求、代码、PRD 放在一起做高密度对照。"),
        ("联动同步提效", "规则拍板后，把 PRD、原型和章节口径更快同步。"),
        ("问题暴露更早", "很多原本靠人工来回翻材料才能发现的问题，会更早出现。"),
    ]
    top = 2.6
    for title, desc in value_items:
        add_box(slide, Inches(0.95), Inches(top), Inches(3.75), Inches(0.78), fill=RGBColor(30, 51, 95), line=None, radius=True)
        add_text(slide, Inches(1.12), Inches(top + 0.08), Inches(3.45), Inches(0.2), title, font_size=14, bold=True, color=WHITE)
        add_text(slide, Inches(1.12), Inches(top + 0.32), Inches(3.35), Inches(0.32), desc, font_size=10, color=RGBColor(215, 223, 243))
        top += 0.95

    add_box(slide, Inches(5.25), Inches(1.7), Inches(7.55), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, radius=True)
    add_text(slide, Inches(5.55), Inches(2.0), Inches(4.0), Inches(0.45), "这 4 类价值，对产品经理意味着什么",
             font_size=18, bold=True, color=INK)

    card_pos = [
        (Inches(5.55), Inches(2.55)),
        (Inches(9.15), Inches(2.55)),
        (Inches(5.55), Inches(4.05)),
        (Inches(9.15), Inches(4.05)),
    ]
    cards = [
        ("01", "少做机械整理", "把时间从搬运信息，转到判断信息。"),
        ("02", "少做重复核对", "高耗时的对照工作，更适合先交给 AI。"),
        ("03", "少被碎活拖住", "文档和原型同步这种碎工作，提效最明显。"),
        ("04", "更聚焦高价值判断", "边界、口径、风险和取舍，仍然必须由产品经理拍板。"),
    ]
    for (left, top), (num, title, desc) in zip(card_pos, cards):
        add_box(slide, left, top, Inches(3.0), Inches(1.1), fill=SAND, line=LIGHT_BLUE, radius=True)
        add_box(slide, left + Inches(0.18), top + Inches(0.14), Inches(0.36), Inches(0.36), num, fill=BLUE, line=None, radius=True,
                font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)
        add_text(slide, left + Inches(0.66), top + Inches(0.14), Inches(2.1), Inches(0.2), title, font_size=13, bold=True, color=INK)
        add_text(slide, left + Inches(0.18), top + Inches(0.56), Inches(2.6), Inches(0.34), desc, font_size=10, color=GRAY)

    add_box(slide, Inches(5.55), Inches(5.58), Inches(6.95), Inches(0.78), fill=SOFT_BLUE, line=BLUE, radius=True)
    add_text(slide, Inches(5.78), Inches(5.82), Inches(6.45), Inches(0.28),
             "AI 的提效是成立的，但这个提效有边界。它节省的是整理、核查、同步和回写时间，不替代业务判断。",
             font_size=11, bold=True, color=INK)

    add_slide_number(slide, 2, 3)


def slide_case(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, fill=RGBColor(238, 243, 251))
    add_text(slide, Inches(0.7), Inches(0.5), Inches(2.2), Inches(0.22), "CHAPTER 4 / CASE 4.5",
             font_size=10, bold=True, color=RED)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(6.4), Inches(0.55),
             "手机端件明细展示方式改造", font_size=24, bold=True, color=INK)
    add_box(slide, Inches(11.25), Inches(0.48), Inches(1.45), Inches(0.34), "案例对比页", fill=PAPER, line=LIGHT_BLUE, radius=True,
            font_size=9, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)

    add_box(slide, Inches(0.7), Inches(1.65), Inches(4.9), Inches(3.95), fill=SOFT_RED, line=RGBColor(245, 203, 198), radius=True)
    add_box(slide, Inches(7.73), Inches(1.65), Inches(5.0), Inches(3.95), fill=SOFT_BLUE, line=LIGHT_BLUE, radius=True)
    add_box(slide, Inches(5.93), Inches(3.05), Inches(1.1), Inches(1.1), "→", fill=BLUE, line=None, radius=True,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)

    add_box(slide, Inches(1.02), Inches(1.96), Inches(0.92), Inches(0.34), "Before", fill=RGBColor(253, 230, 225), line=None, radius=True,
            font_size=9, bold=True, color=RED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)
    add_text(slide, Inches(0.98), Inches(2.38), Inches(2.2), Inches(0.25), "容易带偏的说法", font_size=16, bold=True, color=INK)
    add_box(slide, Inches(1.0), Inches(2.78), Inches(4.3), Inches(0.86),
            "“提纸单和退纸单手机端的件明细展示方式我们要再考虑下……你思考下有什么改进方案。”",
            fill=WHITE, line=RGBColor(235, 222, 218), radius=True,
            font_size=12, bold=True, color=INK)
    add_bullets(slide, Inches(1.05), Inches(3.88), Inches(4.1), Inches(1.45), [
        "给了 AI 很大的设计自由度",
        "没有锁定“必须轻改”",
        "没有限制“必须保留现有结构”",
        "结果是 AI 直接做出一版完整重构稿",
    ], font_size=11, color=INK)

    add_box(slide, Inches(8.03), Inches(1.96), Inches(0.92), Inches(0.34), "After", fill=RGBColor(224, 235, 255), line=None, radius=True,
            font_size=9, bold=True, color=BLUE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02)
    add_text(slide, Inches(8.0), Inches(2.38), Inches(2.1), Inches(0.25), "更稳的说法", font_size=16, bold=True, color=INK)
    add_box(slide, Inches(8.02), Inches(2.78), Inches(4.38), Inches(0.86),
            "“先基于现有页面分析问题，只给轻改方向；如果改动超过原结构，先不要直接改原型。”",
            fill=WHITE, line=LIGHT_BLUE, radius=True,
            font_size=12, bold=True, color=INK)
    add_bullets(slide, Inches(8.08), Inches(3.88), Inches(4.15), Inches(1.45), [
        "先收改造边界，再让 AI 出方案",
        "只改件明细展示，不重做整页结构",
        "先做最小改动试探，再决定是否扩展",
        "把“设计完整性”让位于“迭代风险可控”",
    ], font_size=11, color=INK)

    add_box(slide, Inches(0.72), Inches(5.88), Inches(6.08), Inches(0.82), fill=PAPER, line=LIGHT_BLUE, radius=True)
    add_text(slide, Inches(0.96), Inches(6.08), Inches(1.8), Inches(0.18), "为什么这个案例值得讲", font_size=10, bold=True, color=GRAY)
    add_text(slide, Inches(0.96), Inches(6.3), Inches(5.55), Inches(0.28),
             "它很典型地暴露了 AI 在存量功能迭代里的一个偏好：只要边界不够清楚，它就容易从“轻改”滑向“重做”。",
             font_size=11, color=INK)

    add_box(slide, Inches(6.98), Inches(5.88), Inches(5.72), Inches(0.82), fill=PAPER, line=LIGHT_BLUE, radius=True)
    add_text(slide, Inches(7.23), Inches(6.08), Inches(1.1), Inches(0.18), "最后沉淀", font_size=10, bold=True, color=GRAY)
    add_text(slide, Inches(7.23), Inches(6.3), Inches(5.1), Inches(0.28),
             "好指令的核心不是更复杂，而是阶段更明确、边界更清楚。",
             font_size=11, color=INK)

    add_slide_number(slide, 3, 3)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_value(prs)
    slide_case(prs)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
