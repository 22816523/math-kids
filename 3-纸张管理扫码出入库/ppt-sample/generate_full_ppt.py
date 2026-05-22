# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).with_name("AI参与原有功能迭代升级的产品设计经验_13页正式版.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(238, 243, 251)
PAPER = RGBColor(255, 253, 248)
NAVY = RGBColor(11, 30, 67)
DEEP = RGBColor(19, 41, 79)
BLUE = RGBColor(46, 103, 248)
LIGHT_BLUE = RGBColor(220, 232, 255)
SOFT_BLUE = RGBColor(239, 245, 255)
SOFT_RED = RGBColor(255, 244, 239)
RED = RGBColor(239, 79, 69)
SAND = RGBColor(247, 243, 234)
TEXT = RGBColor(28, 46, 84)
SUB = RGBColor(97, 111, 137)
WHITE = RGBColor(255, 255, 255)

FONT = "Microsoft YaHei"


def add_shape(slide, left, top, width, height, *, fill=None, line=None, rounded=False):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=14, color=TEXT, gap=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(gap)
    return box


def add_chip(slide, left, top, text, *, fill=PAPER, line=LIGHT_BLUE, color=TEXT, width=1.75):
    add_shape(slide, left, top, width, Inches(0.34), fill=fill, line=line, rounded=True)
    add_text(
        slide,
        left,
        top + Inches(0.02),
        width,
        Inches(0.22),
        text,
        font_size=10,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
    )


def add_footer(slide, number, total):
    add_shape(
        slide,
        Inches(0.4),
        Inches(6.88),
        Inches(1.0),
        Inches(0.28),
        fill=PAPER,
        line=LIGHT_BLUE,
        rounded=True,
    )
    add_text(
        slide,
        Inches(0.4),
        Inches(6.95),
        Inches(1.0),
        Inches(0.12),
        f"{number:02d} / {total:02d}",
        font_size=9,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )


def add_base(slide, number, total):
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    add_footer(slide, number, total)


def add_title(slide, chapter, title, number, total, *, tag="企业复盘型 / Slidesgo 风格"):
    add_base(slide, number, total)
    add_chip(slide, Inches(0.68), Inches(0.42), chapter, width=1.55)
    add_text(slide, Inches(0.7), Inches(0.95), Inches(7.3), Inches(0.6), title, font_size=24, bold=True)
    add_chip(slide, Inches(10.7), Inches(0.42), tag, width=2.2)


def cover(slide, total):
    add_base(slide, 1, total)
    add_shape(slide, Inches(0.55), Inches(0.52), Inches(7.35), Inches(5.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_shape(slide, Inches(8.1), Inches(0.52), Inches(4.68), Inches(3.78), fill=NAVY, rounded=True)
    add_shape(slide, Inches(8.1), Inches(4.48), Inches(4.68), Inches(2.0), fill=SOFT_BLUE, line=LIGHT_BLUE, rounded=True)

    add_chip(slide, Inches(0.92), Inches(0.92), "AI + PRODUCT", width=1.45)
    add_text(slide, Inches(0.96), Inches(1.6), Inches(6.5), Inches(1.2), "AI参与原有功能迭代升级的\n产品设计经验", font_size=26, bold=True)
    add_text(
        slide,
        Inches(0.96),
        Inches(3.2),
        Inches(6.2),
        Inches(0.7),
        "以纸张管理扫码出入库项目为例，讲清 AI 在存量 ERP 功能迭代里如何给产品经理提效，以及如何避免失控。",
        font_size=14,
        color=SUB,
    )

    chips = [
        ("原有 ERP 迭代", Inches(0.96)),
        ("产品方法复盘", Inches(2.38)),
        ("AI 协作提效", Inches(3.86)),
    ]
    for text, left in chips:
        add_chip(slide, left, Inches(4.42), text, fill=WHITE, line=LIGHT_BLUE, width=1.2)

    add_text(slide, Inches(8.45), Inches(1.0), Inches(2.0), Inches(0.25), "核心判断", font_size=11, bold=True, color=LIGHT_BLUE)
    add_text(
        slide,
        Inches(8.45),
        Inches(1.72),
        Inches(3.6),
        Inches(1.4),
        "AI 的价值，不在于替产品经理做判断，而在于更快地整理、对照、同步和暴露问题。",
        font_size=18,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        Inches(8.45),
        Inches(3.45),
        Inches(3.75),
        Inches(0.28),
        "纸张管理扫码出入库 / 内部经验分享",
        font_size=10,
        color=RGBColor(205, 215, 239),
    )

    stats = [
        ("1", "个存量 ERP\n迭代案例"),
        ("7", "章分享主线"),
        ("3+3", "踩坑与有效\n做法"),
    ]
    lefts = [Inches(8.35), Inches(9.93), Inches(11.51)]
    for left, (num, label) in zip(lefts, stats):
        add_shape(slide, left, Inches(4.9), Inches(1.25), Inches(1.2), fill=WHITE, line=LIGHT_BLUE, rounded=True)
        add_text(slide, left, Inches(5.02), Inches(1.25), Inches(0.25), num, font_size=18, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.05), Inches(5.45), Inches(1.15), Inches(0.35), label, font_size=8, color=SUB, align=PP_ALIGN.CENTER)


def slide_why(slide, total):
    add_title(slide, "CHAPTER 1", "为什么讲这个项目", 2, total)
    add_shape(slide, Inches(0.72), Inches(1.68), Inches(5.15), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(1.02), Inches(1.98), Inches(4.2), Inches(0.35), "为什么适合拿来分享", font_size=18, bold=True)
    add_bullets(
        slide,
        Inches(1.02),
        Inches(2.5),
        Inches(4.4),
        Inches(3.4),
        [
            "不是从 0 到 1，而是原有 ERP 功能迭代升级",
            "更接近大多数产品经理真实工作场景",
            "原始需求、Codex 对话、PRD、原型都有留痕",
            "能完整复盘 AI 介入过程，而不是只看结果",
        ],
        font_size=15,
    )
    add_shape(slide, Inches(6.1), Inches(1.68), Inches(6.55), Inches(4.95), fill=NAVY, rounded=True)
    add_text(slide, Inches(6.45), Inches(2.0), Inches(1.8), Inches(0.25), "项目背景", font_size=12, bold=True, color=LIGHT_BLUE)
    add_text(
        slide,
        Inches(6.45),
        Inches(2.38),
        Inches(5.6),
        Inches(1.65),
        "在已有纸张管理能力上，从“纸张品类/规格管理”升级为支持一物一码、扫码出入库、库位管理的件级管理能力。",
        font_size=20,
        bold=True,
        color=WHITE,
    )
    add_bullets(
        slide,
        Inches(6.45),
        Inches(4.35),
        Inches(5.3),
        Inches(1.7),
        [
            "核心模块：随货同行单、提纸单、退纸单、一物一码查询",
            "关键议题：库位接入、二维码流转、历史库存初始化",
            "难点不是想不出方案，而是怎么不脱离现有系统现实",
        ],
        font_size=13,
        color=WHITE,
    )


def slide_pain(slide, total):
    add_title(slide, "CHAPTER 2", "没有 AI 时，产品经理卡在哪里", 3, total)
    add_shape(slide, Inches(0.72), Inches(1.7), Inches(12.0), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(0.98), Inches(2.0), Inches(4.2), Inches(0.3), "原有功能迭代的典型卡点", font_size=18, bold=True)
    cards = [
        ("需求零散", "信息跨多轮对话分散，前后补充多，容易漏。"),
        ("材料打架", "原始需求、代码、PRD 不在一个口径上。"),
        ("联动很碎", "原型、文档、规则要同步回写，机械活多。"),
        ("范围失控", "一不留神就从轻改滑向重做，越做越大。"),
    ]
    positions = [
        (Inches(1.0), Inches(2.55)),
        (Inches(4.0), Inches(2.55)),
        (Inches(7.0), Inches(2.55)),
        (Inches(10.0), Inches(2.55)),
    ]
    for (left, top), (title, desc) in zip(positions, cards):
        add_shape(slide, left, top, Inches(2.5), Inches(2.25), fill=SOFT_BLUE, line=LIGHT_BLUE, rounded=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.18), Inches(2.0), Inches(0.25), title, font_size=16, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.62), Inches(2.05), Inches(1.1), desc, font_size=12, color=SUB)
    add_shape(slide, Inches(1.0), Inches(5.2), Inches(11.0), Inches(0.9), fill=SAND, line=LIGHT_BLUE, rounded=True)
    add_text(
        slide,
        Inches(1.22),
        Inches(5.5),
        Inches(10.4),
        Inches(0.26),
        "这个项目真正难的，不是“想不出方案”，而是“怎么贴着现有系统现实收口”。",
        font_size=14,
        bold=True,
    )


def slide_ai_steps(slide, total):
    add_title(slide, "CHAPTER 3", "AI 是怎么介入的", 4, total)
    add_shape(slide, Inches(0.72), Inches(1.7), Inches(12.0), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(0.98), Inches(2.0), Inches(3.6), Inches(0.3), "AI 主要介入的四个环节", font_size=18, bold=True)
    steps = [
        ("01", "需求整理", "把零散原始需求先收成结构稿"),
        ("02", "三方对照", "原始需求 + 代码 + PRD 高密度核查"),
        ("03", "联动同步", "规则拍板后同步文档与原型口径"),
        ("04", "问题暴露", "更早暴露冲突、遗漏和不一致"),
    ]
    left = Inches(1.0)
    for idx, (num, title, desc) in enumerate(steps):
        add_shape(slide, left, Inches(2.75), Inches(2.45), Inches(2.35), fill=SOFT_BLUE if idx % 2 == 0 else SAND, line=LIGHT_BLUE, rounded=True)
        add_shape(slide, left + Inches(0.18), Inches(2.92), Inches(0.42), Inches(0.42), fill=BLUE, rounded=True)
        add_text(slide, left + Inches(0.18), Inches(3.02), Inches(0.42), Inches(0.15), num, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.72), Inches(2.96), Inches(1.5), Inches(0.2), title, font_size=15, bold=True)
        add_text(slide, left + Inches(0.18), Inches(3.55), Inches(1.95), Inches(0.62), desc, font_size=11, color=SUB)
        if idx < 3:
            add_text(slide, left + Inches(2.5), Inches(3.48), Inches(0.35), Inches(0.18), "→", font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        left += Inches(2.95)
    add_shape(slide, Inches(1.0), Inches(5.35), Inches(10.9), Inches(0.8), fill=NAVY, rounded=True)
    add_text(
        slide,
        Inches(1.25),
        Inches(5.6),
        Inches(10.4),
        Inches(0.22),
        "这里最值钱的不是“让 AI 写”，而是让 AI 先对、先找、先暴露，再由产品经理拍板。",
        font_size=13,
        bold=True,
        color=WHITE,
    )


def slide_ai_value(slide, total):
    add_title(slide, "CHAPTER 3", "AI 的价值体现在哪", 5, total)
    add_shape(slide, Inches(0.72), Inches(1.7), Inches(4.45), Inches(4.95), fill=NAVY, rounded=True)
    add_text(slide, Inches(1.0), Inches(2.0), Inches(3.7), Inches(0.35), "AI 对产品经理的提效", font_size=18, bold=True, color=WHITE)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.58),
        Inches(3.6),
        Inches(2.9),
        [
            "原始需求整理更快",
            "差异核查更快",
            "联动同步更快",
            "问题暴露更早",
        ],
        font_size=15,
        color=WHITE,
    )
    add_text(slide, Inches(1.0), Inches(5.65), Inches(3.5), Inches(0.5), "不是少思考了，而是少做机械整理和重复核对。", font_size=12, color=LIGHT_BLUE)

    add_shape(slide, Inches(5.45), Inches(1.7), Inches(7.27), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(5.75), Inches(2.0), Inches(4.0), Inches(0.35), "提效成立，但边界也很清楚", font_size=18, bold=True)
    cards = [
        ("节省的时间", "整理、核查、同步、回写"),
        ("保留给 PM 的事", "边界判断、规则拍板、最终取舍"),
        ("AI 不擅长", "自动理解真实业务边界"),
        ("因此最优模式", "人拍板，AI 做大部分高耗时协作"),
    ]
    positions = [
        (Inches(5.78), Inches(2.68)),
        (Inches(9.12), Inches(2.68)),
        (Inches(5.78), Inches(4.15)),
        (Inches(9.12), Inches(4.15)),
    ]
    for (left, top), (title, desc) in zip(positions, cards):
        add_shape(slide, left, top, Inches(2.9), Inches(1.08), fill=SAND, line=LIGHT_BLUE, rounded=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.14), Inches(2.35), Inches(0.2), title, font_size=13, bold=True)
        add_text(slide, left + Inches(0.18), top + Inches(0.5), Inches(2.45), Inches(0.26), desc, font_size=11, color=SUB)


def add_case_card(slide, left, top, width, height, title, issue, fix, learn, *, fill):
    add_shape(slide, left, top, width, height, fill=fill, line=LIGHT_BLUE, rounded=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.16), width - Inches(0.36), Inches(0.24), title, font_size=15, bold=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.56), Inches(0.55), Inches(0.18), "问题", font_size=10, bold=True, color=SUB)
    add_text(slide, left + Inches(0.78), top + Inches(0.56), width - Inches(1.0), Inches(0.42), issue, font_size=11, color=TEXT)
    add_text(slide, left + Inches(0.18), top + Inches(1.1), Inches(0.55), Inches(0.18), "纠偏", font_size=10, bold=True, color=SUB)
    add_text(slide, left + Inches(0.78), top + Inches(1.1), width - Inches(1.0), Inches(0.42), fix, font_size=11, color=TEXT)
    add_text(slide, left + Inches(0.18), top + Inches(1.64), Inches(0.55), Inches(0.18), "沉淀", font_size=10, bold=True, color=SUB)
    add_text(slide, left + Inches(0.78), top + Inches(1.64), width - Inches(1.0), Inches(0.52), learn, font_size=11, color=TEXT)


def slide_pitfall_1(slide, total):
    add_title(slide, "CHAPTER 4", "踩坑案例：AI 为什么会失控", 6, total, tag="案例 1-2")
    add_case_card(
        slide,
        Inches(0.72),
        Inches(1.72),
        Inches(5.8),
        Inches(4.95),
        "案例1：误判二维码能力可复用",
        "把产成品二维码当成纸张二维码，默认应复用旧能力。",
        "明确指出原来的码是给产成品用的，不是给纸张用的。",
        "系统里有相似能力，不等于业务上真可复用。",
        fill=SOFT_RED,
    )
    add_case_card(
        slide,
        Inches(6.87),
        Inches(1.72),
        Inches(5.85),
        Inches(4.95),
        "案例2：把 WMS 和初始化带进方案",
        "把实施项和后端处理项默认写成产品功能，范围做大。",
        "明确纸张只在 ERP 内闭环，初始化只做后端导入。",
        "涉及导入、同步、历史数据时，先分清是功能、实施项还是后端处理项。",
        fill=SOFT_BLUE,
    )


def slide_pitfall_2(slide, total):
    add_title(slide, "CHAPTER 4", "踩坑案例：AI 为什么会失控", 7, total, tag="案例 3-4")
    add_case_card(
        slide,
        Inches(0.72),
        Inches(1.72),
        Inches(5.8),
        Inches(4.95),
        "案例3：以为流程不变，改造量其实不小",
        "“流程不变，只加扫码”会低估状态、按钮和数据承载改造量。",
        "回到代码现实，按真实状态机和节点重写 PRD。",
        "业务表达上的“流程不变”，不能直接等于技术工作量很小。",
        fill=SOFT_BLUE,
    )
    add_case_card(
        slide,
        Inches(6.87),
        Inches(1.72),
        Inches(5.85),
        Inches(4.95),
        "案例4：手机端原型一次改太大",
        "AI 按更理想的新方案重构页面，已经接近重新设计。",
        "保留对比版，回到轻改路径，只改件明细展示方式。",
        "在存量系统迭代里，控制改造风险往往比设计更完整更重要。",
        fill=SOFT_RED,
    )


def slide_pitfall_common(slide, total):
    add_title(slide, "CHAPTER 4", "这一类坑的共性", 8, total, tag="共性总结")
    add_shape(slide, Inches(0.72), Inches(1.72), Inches(12.0), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    items = [
        ("套旧经验", "看到相似能力，就默认应该复用。"),
        ("自动扩范围", "把实施项、导入、对接都产品化。"),
        ("理想化重设计", "更愿意给新方案，而不是贴着现有系统轻改。"),
        ("局部对，全局乱", "能完成当前指令，但不会天然维护整份文档一致性。"),
    ]
    positions = [
        (Inches(1.0), Inches(2.42)),
        (Inches(6.7), Inches(2.42)),
        (Inches(1.0), Inches(4.2)),
        (Inches(6.7), Inches(4.2)),
    ]
    for (left, top), (title, desc) in zip(positions, items):
        add_shape(slide, left, top, Inches(4.9), Inches(1.22), fill=SAND if top < Inches(4) else SOFT_BLUE, line=LIGHT_BLUE, rounded=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.18), Inches(1.5), Inches(0.2), title, font_size=15, bold=True)
        add_text(slide, left + Inches(0.2), top + Inches(0.56), Inches(4.2), Inches(0.26), desc, font_size=11, color=SUB)
    add_shape(slide, Inches(1.0), Inches(5.72), Inches(10.75), Inches(0.55), fill=NAVY, rounded=True)
    add_text(
        slide,
        Inches(1.22),
        Inches(5.88),
        Inches(10.3),
        Inches(0.18),
        "共性不是 AI 不会写，而是 AI 会默认把“更像好方案的东西”先写出来。",
        font_size=12,
        bold=True,
        color=WHITE,
    )


def slide_success_1(slide, total):
    add_title(slide, "CHAPTER 5", "成功案例：怎样用 AI 才真正有效", 9, total, tag="正向做法 1")
    add_shape(slide, Inches(0.72), Inches(1.72), Inches(12.0), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(0.98), Inches(2.0), Inches(3.5), Inches(0.25), "先把协作顺序做对", font_size=18, bold=True)
    steps = [
        ("先记录，不分析", "把 AI 锁成记录员，避免过早脑补方案。"),
        ("先三方对照，不直接改", "原始需求、代码、PRD 一起核，先暴露差异。"),
        ("先讲改法，再动正式文档", "先对齐修改策略，再落 PRD 和原型。"),
    ]
    lefts = [Inches(1.0), Inches(4.35), Inches(7.7)]
    for left, (title, desc) in zip(lefts, steps):
        add_shape(slide, left, Inches(2.7), Inches(2.7), Inches(2.55), fill=SOFT_BLUE, line=LIGHT_BLUE, rounded=True)
        add_text(slide, left + Inches(0.18), Inches(3.0), Inches(2.15), Inches(0.45), title, font_size=17, bold=True)
        add_text(slide, left + Inches(0.18), Inches(3.82), Inches(2.15), Inches(0.85), desc, font_size=12, color=SUB)
    add_text(slide, Inches(3.7), Inches(3.7), Inches(0.3), Inches(0.2), "→", font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.05), Inches(3.7), Inches(0.3), Inches(0.2), "→", font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(1.0), Inches(5.78), Inches(10.95), Inches(0.48), fill=SAND, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(1.22), Inches(5.93), Inches(10.4), Inches(0.16), "关键不是提示词更花，而是阶段顺序更对。", font_size=12, bold=True)


def slide_success_2(slide, total):
    add_title(slide, "CHAPTER 5", "成功案例：怎样用 AI 才真正有效", 10, total, tag="正向做法 2")
    add_shape(slide, Inches(0.72), Inches(1.72), Inches(5.2), Inches(4.95), fill=NAVY, rounded=True)
    add_text(slide, Inches(1.0), Inches(2.0), Inches(3.8), Inches(0.28), "把 AI 放在更稳的位置上", font_size=18, bold=True, color=WHITE)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.55),
        Inches(3.95),
        Inches(2.8),
        [
            "强约束范围，小步迭代",
            "原型与 PRD 联动时，把 AI 当“同步器”",
            "每轮都做范围复核，避免 AI 做多",
        ],
        font_size=15,
        color=WHITE,
    )
    add_text(slide, Inches(1.0), Inches(5.7), Inches(3.85), Inches(0.3), "边界越清楚，AI 越稳定。", font_size=12, color=LIGHT_BLUE)

    add_shape(slide, Inches(6.18), Inches(1.72), Inches(6.54), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(6.48), Inches(2.0), Inches(4.2), Inches(0.3), "这部分真正要传递的东西", font_size=18, bold=True)
    add_shape(slide, Inches(6.48), Inches(2.62), Inches(5.9), Inches(0.82), fill=SOFT_BLUE, line=LIGHT_BLUE, rounded=True)
    add_text(slide, Inches(6.72), Inches(2.88), Inches(5.4), Inches(0.22), "不是“提示词写得漂亮”，而是“人先定边界，AI 再高效执行”。", font_size=13, bold=True)
    add_shape(slide, Inches(6.48), Inches(3.72), Inches(5.9), Inches(1.62), fill=SAND, line=LIGHT_BLUE, rounded=True)
    add_bullets(
        slide,
        Inches(6.72),
        Inches(4.02),
        Inches(5.2),
        Inches(1.0),
        [
            "方向未定时，AI 更适合暴露问题，不适合直接拍板。",
            "口径已定后，AI 最适合做同步、回写和一致性维护。",
        ],
        font_size=12,
    )


def slide_instruction(slide, total):
    add_title(slide, "CHAPTER 6", "协作指令模式", 11, total, tag="高价值协作方式")
    add_shape(slide, Inches(0.72), Inches(1.72), Inches(5.75), Inches(4.95), fill=SOFT_BLUE, line=LIGHT_BLUE, rounded=True)
    add_shape(slide, Inches(6.95), Inches(1.72), Inches(5.77), Inches(4.95), fill=SOFT_RED, line=RGBColor(245, 203, 198), rounded=True)
    add_text(slide, Inches(0.98), Inches(2.0), Inches(2.2), Inches(0.25), "高价值说法", font_size=18, bold=True)
    add_text(slide, Inches(7.22), Inches(2.0), Inches(2.6), Inches(0.25), "容易带偏的说法", font_size=18, bold=True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.55),
        Inches(5.0),
        Inches(3.2),
        [
            "在我没有结束之前，你不用做任何分析，只接收和记录。",
            "结合代码和原始需求，对照 PRD 找问题。",
            "先不改，先告诉我你准备怎么改。",
            "手机端先不改，我们先改 PC 端，按照原型改。",
        ],
        font_size=13,
    )
    add_bullets(
        slide,
        Inches(7.22),
        Inches(2.55),
        Inches(5.0),
        Inches(3.2),
        [
            "你思考下有什么改进方案。",
            "你先把详情页原型改了，我看看效果。",
            "这类说法的问题不在语气，而在边界和阶段不够清楚。",
        ],
        font_size=13,
    )
    add_shape(slide, Inches(1.0), Inches(5.95), Inches(11.0), Inches(0.4), fill=NAVY, rounded=True)
    add_text(slide, Inches(1.2), Inches(6.06), Inches(10.6), Inches(0.15), "好指令的核心不是更复杂，而是边界更清楚、阶段更明确。", font_size=12, bold=True, color=WHITE)


def slide_conclusion(slide, total):
    add_title(slide, "CHAPTER 7", "最后想传递的结论", 12, total, tag="结论页")
    add_shape(slide, Inches(0.72), Inches(1.72), Inches(12.0), Inches(4.95), fill=PAPER, line=LIGHT_BLUE, rounded=True)
    add_shape(slide, Inches(1.0), Inches(2.0), Inches(10.9), Inches(0.92), fill=NAVY, rounded=True)
    add_text(
        slide,
        Inches(1.3),
        Inches(2.27),
        Inches(10.3),
        Inches(0.25),
        "AI 的价值，不在于替产品经理做产品，而在于帮你更快地整理、比对、同步和暴露问题。",
        font_size=16,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    bullets = [
        "AI 的提效是成立的，但这个提效有边界。",
        "真正决定产出质量的，仍然是产品经理有没有讲清边界、拍板规则、持续联动校正。",
        "AI 节省的是整理、核查、同步和回写时间，不替代业务判断。",
        "最值得沉淀的不是一组提示词，而是一套适用于存量功能迭代的协作方法。",
    ]
    add_bullets(slide, Inches(1.18), Inches(3.42), Inches(10.5), Inches(2.15), bullets, font_size=14)


def slide_closing(slide, total):
    add_base(slide, 13, total)
    add_shape(slide, Inches(0.95), Inches(0.92), Inches(11.45), Inches(5.7), fill=NAVY, rounded=True)
    add_chip(slide, Inches(1.25), Inches(1.28), "CLOSING", fill=WHITE, line=None, color=DEEP, width=1.2)
    add_text(
        slide,
        Inches(1.28),
        Inches(2.0),
        Inches(10.8),
        Inches(1.55),
        "在原有功能迭代项目里，\nAI 不是替代产品经理的设计者，\n而是放大产品经理能力的协作工具。",
        font_size=26,
        bold=True,
        color=WHITE,
    )
    add_bullets(
        slide,
        Inches(1.32),
        Inches(4.45),
        Inches(10.5),
        Inches(1.2),
        [
            "先定边界，再让 AI 展开",
            "先对照现实，再让 AI 输出",
            "先拍板口径，再让 AI 同步",
            "先控制范围，再让 AI 提效",
        ],
        font_size=14,
        color=LIGHT_BLUE,
    )


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        cover,
        slide_why,
        slide_pain,
        slide_ai_steps,
        slide_ai_value,
        slide_pitfall_1,
        slide_pitfall_2,
        slide_pitfall_common,
        slide_success_1,
        slide_success_2,
        slide_instruction,
        slide_conclusion,
        slide_closing,
    ]

    total = len(slides)
    for builder in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder(slide, total)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
