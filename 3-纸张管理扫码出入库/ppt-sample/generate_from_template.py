# -*- coding: utf-8 -*-
from copy import deepcopy
from pathlib import Path

from pptx import Presentation


TEMPLATE = Path(r"D:\文档\JTY\AI项目文档\3-纸张管理扫码出入库\ppt-sample\PPT模板.pptx")
OUTPUT = TEMPLATE.with_name("AI参与原有功能迭代升级的产品设计经验_套模板版.pptx")


def set_text(slide, index, value):
    slide.shapes[index - 1].text = value


def duplicate_slide(prs, index):
    source = prs.slides[index - 1]
    blank = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank)
    for shape in list(new_slide.shapes):
        el = shape._element
        el.getparent().remove(el)
    for shape in source.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")
    return new_slide


def build():
    template = Presentation(str(TEMPLATE))
    out = Presentation(str(TEMPLATE))

    for idx in range(len(out.slides) - 1, -1, -1):
        r_id = out.slides._sldIdLst[idx].rId
        out.part.drop_rel(r_id)
        del out.slides._sldIdLst[idx]

    selected = [1, 2, 5, 6, 17, 4, 10, 13, 18, 19, 20]
    for idx in selected:
        source = template.slides[idx - 1]
        blank = out.slide_layouts[0]
        slide = out.slides.add_slide(blank)
        for shape in list(slide.shapes):
            el = shape._element
            el.getparent().remove(el)
        for shape in source.shapes:
            slide.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")

    # 1 封面
    slide = out.slides[0]
    set_text(slide, 5, "AI参与原有功能迭代升级\n产品设计经验分享")
    set_text(slide, 7, "AI-ENABLED PRODUCT ITERATION")
    set_text(slide, 10, "Paper management scan in/out inventory project case")

    # 2 目录
    slide = out.slides[1]
    set_text(slide, 1, "CONTENTS")
    set_text(slide, 2, "目录-CONTENTS")
    set_text(slide, 6, "为什么讲这个项目")
    set_text(slide, 7, "Why This Project")
    set_text(slide, 9, "01")
    set_text(slide, 16, "AI是怎么介入的，价值体现在哪")
    set_text(slide, 17, "How AI Added Value")
    set_text(slide, 19, "02")
    set_text(slide, 11, "踩坑案例与成功案例")
    set_text(slide, 12, "Pitfalls And Effective Practices")
    set_text(slide, 14, "03")
    set_text(slide, 21, "协作指令模式与最终结论")
    set_text(slide, 22, "Prompting Patterns And Final Takeaways")
    set_text(slide, 24, "04")
    set_text(slide, 26, "AI product design experience sharing")

    # 3 章节页：为什么讲这个项目
    slide = out.slides[2]
    set_text(slide, 1, "为什么讲这个项目")
    set_text(slide, 2, "Why This Project")

    # 4 内容页：AI价值
    slide = out.slides[3]
    set_text(slide, 1, "AI 的价值体现在哪")
    set_text(slide, 2, "Where AI Created Value")
    set_text(slide, 3, "这次项目里，AI 的价值不在于替产品经理做判断，而在于更快地整理、对照、同步和暴露问题。")
    set_text(slide, 4, "价值体现")
    set_text(slide, 5, "提效")

    # 5 三列总结页：无AI卡点
    slide = out.slides[4]
    set_text(slide, 1, "没有 AI 时，产品经理卡在哪里")
    set_text(slide, 2, "Pain Points Without AI")
    set_text(slide, 4, "需求零散")
    set_text(slide, 5, "原始需求跨多轮对话分散，前后补充多，容易漏。")
    set_text(slide, 6, "01")
    set_text(slide, 9, "材料打架")
    set_text(slide, 10, "原始需求、代码、PRD 经常不在一个口径上。")
    set_text(slide, 8, "02")
    set_text(slide, 13, "联动很碎")
    set_text(slide, 14, "原型、文档、规则要同步回写，机械活很多。")
    set_text(slide, 12, "03")

    # 6 五点特征页：为什么讲这个项目
    slide = out.slides[5]
    set_text(slide, 1, "为什么选这个项目")
    set_text(slide, 2, "Why This Case")
    set_text(slide, 9, "不是从0到1，而是原有ERP迭代升级")
    set_text(slide, 11, "更接近大多数产品经理的真实工作场景")
    set_text(slide, 12, "原始需求、Codex对话、PRD、原型都有留痕")
    set_text(slide, 13, "能完整复盘 AI 的价值和边界")
    set_text(slide, 14, "难点不是想不出方案，而是怎么贴着现有系统收口")
    set_text(slide, 25, "五个判断")

    # 7 四点卡片：AI怎么介入
    slide = out.slides[6]
    set_text(slide, 1, "AI 是怎么介入的")
    set_text(slide, 2, "How AI Was Used")
    set_text(slide, 11, "需求整理")
    set_text(slide, 12, "先把零散需求收成结构草稿。")
    set_text(slide, 13, "三方对照")
    set_text(slide, 14, "把原始需求、代码、PRD 放在一起高密度核查。")
    set_text(slide, 15, "联动同步")
    set_text(slide, 16, "规则拍板后，同步回写文档和原型。")
    set_text(slide, 17, "问题暴露")
    set_text(slide, 18, "更早暴露冲突、遗漏和不一致。")

    # 8 四象限页：踩坑共性
    slide = out.slides[7]
    set_text(slide, 1, "踩坑案例：AI 为什么会失控")
    set_text(slide, 2, "Why AI Went Off Track")
    set_text(slide, 7, "套旧经验")
    set_text(slide, 8, "看到相似能力，就默认应该复用。")
    set_text(slide, 10, "自动扩范围")
    set_text(slide, 11, "把实施项、导入、对接都产品化。")
    set_text(slide, 13, "理想化重设计")
    set_text(slide, 14, "更愿意给新方案，而不是贴着现有系统轻改。")
    set_text(slide, 16, "局部对，全局乱")
    set_text(slide, 17, "能完成当前指令，但不会天然维护整份文档一致性。")

    # 9 三块内容：成功做法
    slide = out.slides[8]
    set_text(slide, 1, "成功案例：怎样用 AI 才真正有效")
    set_text(slide, 2, "What Actually Worked")
    set_text(slide, 17, "先记录，不分析")
    set_text(slide, 18, "先把 AI 锁成记录员，避免过早脑补方案。")
    set_text(slide, 19, "先三方对照")
    set_text(slide, 20, "原始需求、代码、PRD 一起核，先暴露差异。")
    set_text(slide, 21, "范围压小，再让 AI 展开")
    set_text(slide, 22, "边界越清楚，AI 越稳定。")

    # 10 多标签页：协作指令模式
    slide = out.slides[9]
    set_text(slide, 1, "协作指令模式")
    set_text(slide, 2, "Prompting Patterns")
    set_text(slide, 5, "实际有效")
    set_text(slide, 6, "在我没有结束之前，你不用做任何分析，只接收和记录。")
    set_text(slide, 7, "实际有效")
    set_text(slide, 8, "结合代码和原始需求，对照 PRD 找问题。")
    set_text(slide, 9, "实际有效")
    set_text(slide, 10, "先不改，先告诉我你准备怎么改。")
    set_text(slide, 11, "实际有效")
    set_text(slide, 12, "手机端先不改，我们先改 PC 端，按照原型改。")
    set_text(slide, 13, "容易带偏")
    set_text(slide, 14, "你思考下有什么改进方案。")
    set_text(slide, 15, "容易带偏")
    set_text(slide, 16, "你先把详情页原型改了，我看看效果。")
    set_text(slide, 17, "一句话")
    set_text(slide, 18, "好指令的核心不是更复杂，而是边界更清楚、阶段更明确。")
    set_text(slide, 19, "总结")

    # 11 结束结论页
    slide = out.slides[10]
    set_text(slide, 5, "最后想传递的结论\nTHANK YOU.")
    set_text(slide, 7, "AI-ENABLED PRODUCT ITERATION")
    set_text(slide, 10, "AI 的价值，不在于替产品经理做产品，而在于帮你更快地整理、比对、同步和暴露问题。")

    out.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
