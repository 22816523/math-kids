# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).with_name("AI参与原有功能迭代升级的产品设计经验_14页美化版.pptx")

SW = Inches(13.333)
SH = Inches(7.5)

FONT = "Microsoft YaHei"

NAVY = RGBColor(15, 33, 72)
INDIGO = RGBColor(33, 71, 158)
BLUE = RGBColor(53, 110, 255)
CYAN = RGBColor(103, 187, 255)
BG = RGBColor(244, 247, 252)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(28, 46, 84)
SUB = RGBColor(95, 111, 142)
LINE = RGBColor(220, 230, 247)
SOFT_BLUE = RGBColor(234, 242, 255)
SOFT_CYAN = RGBColor(236, 248, 255)
SOFT_RED = RGBColor(255, 242, 239)
SOFT_SAND = RGBColor(248, 245, 238)
RED = RGBColor(234, 88, 78)
GOLD = RGBColor(255, 186, 73)


def rect(slide, left, top, width, height, *, fill=None, line=None, rounded=False, oval=False):
    if oval:
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
    else:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if hasattr(shape.fill, "transparency"):
            shape.fill.transparency = 0
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def text(
    slide,
    left,
    top,
    width,
    height,
    value,
    *,
    size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, left, top, width, height, items, *, size=14, color=TEXT, gap=7):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(gap)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def footer(slide, number, total):
    rect(slide, Inches(11.75), Inches(6.82), Inches(1.15), Inches(0.34), fill=WHITE, line=LINE, rounded=True)
    text(slide, Inches(11.75), Inches(6.91), Inches(1.15), Inches(0.12), f"{number:02d}/{total:02d}", size=9, bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def chip(slide, left, top, width, label, *, fill=WHITE, line=LINE, color=TEXT):
    rect(slide, left, top, width, Inches(0.34), fill=fill, line=line, rounded=True)
    text(slide, left, top + Inches(0.08), width, Inches(0.12), label, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)


def base_bg(slide, number, total, *, accent=BLUE, chapter=None, title=None):
    rect(slide, 0, 0, SW, SH, fill=BG)
    rect(slide, 0, 0, SW, Inches(0.2), fill=accent)
    rect(slide, Inches(9.8), Inches(-0.9), Inches(4.4), Inches(2.8), fill=SOFT_BLUE, oval=True)
    rect(slide, Inches(10.8), Inches(-0.5), Inches(2.2), Inches(1.4), fill=SOFT_CYAN, oval=True)
    rect(slide, Inches(-0.7), Inches(5.85), Inches(2.4), Inches(1.5), fill=SOFT_SAND, oval=True)
    rect(slide, Inches(0.72), Inches(0.62), Inches(12.0), Inches(5.95), fill=WHITE, line=LINE, rounded=True)
    rect(slide, Inches(0.72), Inches(0.62), Inches(0.22), Inches(5.95), fill=accent, rounded=True)
    if chapter:
        chip(slide, Inches(1.08), Inches(0.95), Inches(1.45), chapter, fill=SOFT_BLUE, line=None, color=accent)
    if title:
        text(slide, Inches(1.08), Inches(1.48), Inches(7.8), Inches(0.42), title, size=24, bold=True)
    footer(slide, number, total)


def cover(slide, total):
    rect(slide, 0, 0, SW, SH, fill=NAVY)
    rect(slide, Inches(8.7), Inches(-0.8), Inches(5.2), Inches(3.0), fill=INDIGO, oval=True)
    rect(slide, Inches(9.8), Inches(0.2), Inches(3.0), Inches(2.0), fill=BLUE, oval=True)
    rect(slide, Inches(-0.5), Inches(5.6), Inches(3.2), Inches(2.0), fill=INDIGO, oval=True)
    rect(slide, Inches(0.78), Inches(0.78), Inches(11.7), Inches(5.9), fill=RGBColor(23, 45, 92), line=RGBColor(65, 91, 151), rounded=True)
    rect(slide, Inches(0.96), Inches(1.0), Inches(6.9), Inches(5.45), fill=WHITE, line=None, rounded=True)
    rect(slide, Inches(8.08), Inches(1.0), Inches(3.95), Inches(3.55), fill=BLUE, rounded=True)
    rect(slide, Inches(8.08), Inches(4.75), Inches(3.95), Inches(1.7), fill=SOFT_CYAN, line=None, rounded=True)

    chip(slide, Inches(1.22), Inches(1.3), Inches(1.35), "AI + PRODUCT", fill=SOFT_BLUE, line=None, color=BLUE)
    text(slide, Inches(1.22), Inches(2.0), Inches(5.95), Inches(1.5), "AI参与原有功能迭代升级的\n产品设计经验", size=28, bold=True)
    text(
        slide,
        Inches(1.22),
        Inches(3.72),
        Inches(5.6),
        Inches(0.78),
        "以纸张管理扫码出入库项目为例，复盘 AI 在存量 ERP 功能迭代里如何真正给产品经理提效。",
        size=15,
        color=SUB,
    )
    chip(slide, Inches(1.22), Inches(4.7), Inches(1.22), "原有ERP迭代", fill=WHITE, line=LINE)
    chip(slide, Inches(2.62), Inches(4.7), Inches(1.22), "案例复盘", fill=WHITE, line=LINE)
    chip(slide, Inches(4.02), Inches(4.7), Inches(1.22), "协作提效", fill=WHITE, line=LINE)

    text(slide, Inches(8.38), Inches(1.4), Inches(1.5), Inches(0.25), "核心判断", size=12, bold=True, color=WHITE)
    text(
        slide,
        Inches(8.38),
        Inches(1.95),
        Inches(3.15),
        Inches(1.6),
        "AI 的价值，不在于替产品经理做判断，而在于更快地整理、对照、同步和暴露问题。",
        size=19,
        bold=True,
        color=WHITE,
    )
    stats = [("1", "个存量ERP\n案例"), ("7", "章分享主线"), ("3+3", "踩坑与有效\n做法")]
    lefts = [Inches(8.33), Inches(9.67), Inches(11.01)]
    for left, (num, label) in zip(lefts, stats):
        rect(slide, left, Inches(4.98), Inches(1.05), Inches(1.06), fill=WHITE, rounded=True)
        text(slide, left, Inches(5.1), Inches(1.05), Inches(0.22), num, size=18, bold=True, align=PP_ALIGN.CENTER)
        text(slide, left + Inches(0.02), Inches(5.46), Inches(1.0), Inches(0.3), label, size=8, color=SUB, align=PP_ALIGN.CENTER)
    footer(slide, 1, total)


def agenda(slide, total):
    base_bg(slide, 2, total, accent=CYAN, chapter="AGENDA", title="目录")
    rect(slide, Inches(1.12), Inches(2.05), Inches(2.1), Inches(3.95), fill=NAVY, rounded=True)
    text(slide, Inches(1.42), Inches(2.55), Inches(1.5), Inches(1.2), "这次分享\n怎么讲", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(slide, Inches(1.38), Inches(4.65), Inches(1.6), Inches(0.5), "7 个部分", size=12, color=RGBColor(204, 217, 245), align=PP_ALIGN.CENTER)
    items = [
        "为什么讲这个项目",
        "没有 AI 时，产品经理卡在哪里",
        "AI 是怎么介入的，价值体现在哪",
        "踩坑案例：AI 为什么会失控",
        "成功案例：怎样用 AI 才真正有效",
        "协作指令模式",
        "最后想传递的结论",
    ]
    y = 2.0
    for idx, item in enumerate(items, start=1):
        fill = SOFT_BLUE if idx % 2 else SOFT_SAND
        rect(slide, Inches(3.6), Inches(y), Inches(8.45), Inches(0.56), fill=fill, line=LINE, rounded=True)
        rect(slide, Inches(3.8), Inches(y + 0.1), Inches(0.38), Inches(0.36), fill=BLUE, rounded=True)
        text(slide, Inches(3.8), Inches(y + 0.18), Inches(0.38), Inches(0.1), str(idx), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(slide, Inches(4.35), Inches(y + 0.16), Inches(6.95), Inches(0.14), item, size=14, bold=True if idx in (1, 4, 7) else False)
        y += 0.66


def why(slide, total):
    base_bg(slide, 3, total, accent=BLUE, chapter="CHAPTER 1", title="为什么讲这个项目")
    rect(slide, Inches(1.12), Inches(2.0), Inches(4.55), Inches(3.95), fill=SOFT_BLUE, line=LINE, rounded=True)
    rect(slide, Inches(5.92), Inches(2.0), Inches(5.9), Inches(3.95), fill=NAVY, rounded=True)
    text(slide, Inches(1.36), Inches(2.28), Inches(3.2), Inches(0.25), "为什么适合拿来分享", size=18, bold=True)
    bullets(
        slide,
        Inches(1.36),
        Inches(2.82),
        Inches(3.9),
        Inches(2.55),
        [
            "不是从 0 到 1，而是原有 ERP 功能迭代升级",
            "更接近大多数产品经理真实工作场景",
            "原始需求、Codex 对话、PRD、原型都有留痕",
            "能完整复盘 AI 的价值和边界",
        ],
        size=14,
    )
    text(slide, Inches(6.25), Inches(2.25), Inches(1.8), Inches(0.2), "项目背景", size=12, bold=True, color=CYAN)
    text(
        slide,
        Inches(6.25),
        Inches(2.72),
        Inches(4.9),
        Inches(1.35),
        "在已有纸张管理能力上，从“纸张品类/规格管理”升级为支持一物一码、扫码出入库、库位管理的件级管理能力。",
        size=20,
        bold=True,
        color=WHITE,
    )
    bullets(
        slide,
        Inches(6.25),
        Inches(4.5),
        Inches(4.9),
        Inches(1.2),
        [
            "核心模块：随货同行单、提纸单、退纸单、一物一码查询",
            "关键议题：库位接入、二维码流转、历史库存初始化",
        ],
        size=12,
        color=WHITE,
    )


def pain(slide, total):
    base_bg(slide, 4, total, accent=BLUE, chapter="CHAPTER 2", title="没有 AI 时，产品经理卡在哪里")
    titles = ["需求零散", "材料打架", "联动很碎", "范围失控"]
    descs = [
        "信息跨多轮对话分散，前后补充多，容易漏。",
        "原始需求、代码、PRD 经常不在一个口径上。",
        "原型、文档、规则要同步回写，机械活很多。",
        "一不留神就从轻改滑向重做，越做越大。",
    ]
    xs = [1.12, 4.03, 6.94, 9.85]
    for idx in range(4):
        fill = SOFT_BLUE if idx % 2 == 0 else SOFT_SAND
        rect(slide, Inches(xs[idx]), Inches(2.18), Inches(2.55), Inches(2.9), fill=fill, line=LINE, rounded=True)
        text(slide, Inches(xs[idx] + 0.22), Inches(2.5), Inches(1.7), Inches(0.2), titles[idx], size=16, bold=True)
        text(slide, Inches(xs[idx] + 0.22), Inches(3.1), Inches(2.02), Inches(1.1), descs[idx], size=12, color=SUB)
    rect(slide, Inches(1.12), Inches(5.38), Inches(10.95), Inches(0.62), fill=NAVY, rounded=True)
    text(slide, Inches(1.35), Inches(5.58), Inches(10.4), Inches(0.16), "这个项目真正难的，不是“想不出方案”，而是“怎么贴着现有系统现实收口”。", size=13, bold=True, color=WHITE)


def ai_steps(slide, total):
    base_bg(slide, 5, total, accent=CYAN, chapter="CHAPTER 3", title="AI 是怎么介入的")
    steps = [
        ("01", "需求整理", "把零散原始需求先收成结构稿"),
        ("02", "三方对照", "原始需求 + 代码 + PRD 高密度核查"),
        ("03", "联动同步", "规则拍板后同步文档与原型口径"),
        ("04", "问题暴露", "更早暴露冲突、遗漏和不一致"),
    ]
    xs = [1.2, 4.1, 7.0, 9.9]
    for idx, ((num, title_str, desc), x) in enumerate(zip(steps, xs)):
        rect(slide, Inches(x), Inches(2.28), Inches(2.42), Inches(3.12), fill=SOFT_BLUE if idx % 2 == 0 else SOFT_SAND, line=LINE, rounded=True)
        rect(slide, Inches(x + 0.18), Inches(2.48), Inches(0.45), Inches(0.38), fill=BLUE, rounded=True)
        text(slide, Inches(x + 0.18), Inches(2.57), Inches(0.45), Inches(0.12), num, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(slide, Inches(x + 0.8), Inches(2.52), Inches(1.3), Inches(0.18), title_str, size=15, bold=True)
        text(slide, Inches(x + 0.18), Inches(3.28), Inches(1.9), Inches(1.1), desc, size=12, color=SUB)
        if idx < 3:
            text(slide, Inches(x + 2.48), Inches(3.55), Inches(0.22), Inches(0.2), "→", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rect(slide, Inches(1.2), Inches(5.72), Inches(10.95), Inches(0.45), fill=SOFT_CYAN, line=LINE, rounded=True)
    text(slide, Inches(1.42), Inches(5.86), Inches(10.5), Inches(0.16), "这里最值钱的不是“让 AI 写”，而是让 AI 先对、先找、先暴露。", size=12, bold=True)


def ai_value(slide, total):
    base_bg(slide, 6, total, accent=CYAN, chapter="CHAPTER 3", title="AI 的价值体现在哪")
    rect(slide, Inches(1.12), Inches(2.0), Inches(3.8), Inches(4.0), fill=NAVY, rounded=True)
    text(slide, Inches(1.38), Inches(2.28), Inches(2.7), Inches(0.25), "AI 对产品经理的提效", size=18, bold=True, color=WHITE)
    bullets(
        slide,
        Inches(1.38),
        Inches(2.85),
        Inches(2.9),
        Inches(2.4),
        ["原始需求整理更快", "差异核查更快", "联动同步更快", "问题暴露更早"],
        size=15,
        color=WHITE,
    )
    text(slide, Inches(1.38), Inches(5.42), Inches(2.9), Inches(0.3), "不是少思考了，而是少做机械整理和重复核对。", size=11, color=CYAN)

    info = [
        ("节省的时间", "整理、核查、同步、回写"),
        ("保留给 PM 的事", "边界判断、规则拍板、最终取舍"),
        ("AI 不擅长", "自动理解真实业务边界"),
        ("最优模式", "人拍板，AI 做大部分高耗时协作"),
    ]
    pos = [(5.25, 2.15), (8.55, 2.15), (5.25, 4.1), (8.55, 4.1)]
    for (title_str, desc), (x, y) in zip(info, pos):
        rect(slide, Inches(x), Inches(y), Inches(2.8), Inches(1.48), fill=SOFT_BLUE, line=LINE, rounded=True)
        text(slide, Inches(x + 0.2), Inches(y + 0.22), Inches(2.1), Inches(0.18), title_str, size=13, bold=True)
        text(slide, Inches(x + 0.2), Inches(y + 0.72), Inches(2.2), Inches(0.35), desc, size=11, color=SUB)


def case_card(slide, x, y, w, h, title_str, issue, fix, learn, fill):
    rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=LINE, rounded=True)
    text(slide, Inches(x + 0.22), Inches(y + 0.2), Inches(w - 0.45), Inches(0.25), title_str, size=15, bold=True)
    text(slide, Inches(x + 0.22), Inches(y + 0.62), Inches(0.5), Inches(0.12), "问题", size=10, bold=True, color=SUB)
    text(slide, Inches(x + 0.8), Inches(y + 0.62), Inches(w - 1.05), Inches(0.45), issue, size=11)
    text(slide, Inches(x + 0.22), Inches(y + 1.2), Inches(0.5), Inches(0.12), "纠偏", size=10, bold=True, color=SUB)
    text(slide, Inches(x + 0.8), Inches(y + 1.2), Inches(w - 1.05), Inches(0.45), fix, size=11)
    text(slide, Inches(x + 0.22), Inches(y + 1.78), Inches(0.5), Inches(0.12), "沉淀", size=10, bold=True, color=SUB)
    text(slide, Inches(x + 0.8), Inches(y + 1.78), Inches(w - 1.05), Inches(0.5), learn, size=11)


def pitfall1(slide, total):
    base_bg(slide, 7, total, accent=RED, chapter="CHAPTER 4", title="踩坑案例：AI 为什么会失控")
    case_card(slide, 1.12, 2.0, 5.1, 3.95, "案例1：误判二维码能力可复用", "把产成品二维码当成纸张二维码，默认应复用旧能力。", "明确指出原来的码是给产成品用的，不是给纸张用的。", "系统里有相似能力，不等于业务上真可复用。", SOFT_RED)
    case_card(slide, 6.55, 2.0, 5.27, 3.95, "案例2：把 WMS 和初始化带进方案", "把实施项和后端处理项默认写成产品功能，范围做大。", "明确纸张只在 ERP 内闭环，初始化只做后端导入。", "导入、同步、历史数据要先分清是功能、实施项还是后端处理项。", SOFT_BLUE)


def pitfall2(slide, total):
    base_bg(slide, 8, total, accent=RED, chapter="CHAPTER 4", title="踩坑案例：AI 为什么会失控")
    case_card(slide, 1.12, 2.0, 5.1, 3.95, "案例3：以为流程不变，改造量其实不小", "“流程不变，只加扫码”会低估状态、按钮和数据承载改造量。", "回到代码现实，按真实状态机和节点重写 PRD。", "业务表达上的“流程不变”，不能直接等于技术工作量很小。", SOFT_BLUE)
    case_card(slide, 6.55, 2.0, 5.27, 3.95, "案例4：手机端原型一次改太大", "AI 按更理想的新方案重构页面，已经接近重新设计。", "保留对比版，回到轻改路径，只改件明细展示方式。", "在存量系统迭代里，控制改造风险往往比设计更完整更重要。", SOFT_RED)


def pitfall_common(slide, total):
    base_bg(slide, 9, total, accent=RED, chapter="CHAPTER 4", title="这一类坑的共性")
    cards = [
        ("套旧经验", "看到相似能力，就默认应该复用。"),
        ("自动扩范围", "把实施项、导入、对接都产品化。"),
        ("理想化重设计", "更愿意给新方案，而不是贴着现有系统轻改。"),
        ("局部对，全局乱", "能完成当前指令，但不会天然维护整份文档一致性。"),
    ]
    pos = [(1.12, 2.2), (6.55, 2.2), (1.12, 4.2), (6.55, 4.2)]
    for i, ((title_str, desc), (x, y)) in enumerate(zip(cards, pos)):
        rect(slide, Inches(x), Inches(y), Inches(5.0), Inches(1.35), fill=SOFT_SAND if i % 2 else SOFT_BLUE, line=LINE, rounded=True)
        text(slide, Inches(x + 0.22), Inches(y + 0.22), Inches(1.4), Inches(0.18), title_str, size=15, bold=True)
        text(slide, Inches(x + 0.22), Inches(y + 0.72), Inches(4.2), Inches(0.24), desc, size=11, color=SUB)
    rect(slide, Inches(1.12), Inches(5.95), Inches(10.7), Inches(0.42), fill=NAVY, rounded=True)
    text(slide, Inches(1.34), Inches(6.07), Inches(10.2), Inches(0.14), "共性不是 AI 不会写，而是它会先写出“更像好方案”的东西。", size=12, bold=True, color=WHITE)


def success1(slide, total):
    base_bg(slide, 10, total, accent=BLUE, chapter="CHAPTER 5", title="成功案例：怎样用 AI 才真正有效")
    titles = ["先记录，不分析", "先三方对照，不直接改", "先讲改法，再动正式文档"]
    descs = [
        "把 AI 锁成记录员，避免过早脑补方案。",
        "原始需求、代码、PRD 一起核，先暴露差异。",
        "先对齐修改策略，再落 PRD 和原型。",
    ]
    xs = [1.2, 4.45, 7.7]
    for i in range(3):
        rect(slide, Inches(xs[i]), Inches(2.3), Inches(2.7), Inches(3.05), fill=SOFT_BLUE if i != 1 else SOFT_SAND, line=LINE, rounded=True)
        text(slide, Inches(xs[i] + 0.2), Inches(2.7), Inches(2.0), Inches(0.5), titles[i], size=16, bold=True)
        text(slide, Inches(xs[i] + 0.2), Inches(3.8), Inches(2.0), Inches(0.8), descs[i], size=12, color=SUB)
        if i < 2:
            text(slide, Inches(xs[i] + 2.78), Inches(3.62), Inches(0.25), Inches(0.18), "→", size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rect(slide, Inches(1.2), Inches(5.72), Inches(10.65), Inches(0.42), fill=SOFT_CYAN, line=LINE, rounded=True)
    text(slide, Inches(1.42), Inches(5.85), Inches(10.1), Inches(0.15), "关键不是提示词更花，而是协作顺序更对。", size=12, bold=True)


def success2(slide, total):
    base_bg(slide, 11, total, accent=BLUE, chapter="CHAPTER 5", title="成功案例：怎样用 AI 才真正有效")
    rect(slide, Inches(1.12), Inches(2.0), Inches(4.25), Inches(3.95), fill=NAVY, rounded=True)
    text(slide, Inches(1.38), Inches(2.32), Inches(2.7), Inches(0.3), "把 AI 放在更稳的位置上", size=18, bold=True, color=WHITE)
    bullets(
        slide,
        Inches(1.38),
        Inches(2.92),
        Inches(3.2),
        Inches(2.3),
        ["强约束范围，小步迭代", "原型与 PRD 联动时，把 AI 当“同步器”", "每轮都做范围复核，避免 AI 做多"],
        size=14,
        color=WHITE,
    )
    rect(slide, Inches(5.72), Inches(2.0), Inches(6.1), Inches(3.95), fill=SOFT_BLUE, line=LINE, rounded=True)
    text(slide, Inches(6.0), Inches(2.32), Inches(3.5), Inches(0.25), "这部分真正要传递的东西", size=18, bold=True)
    rect(slide, Inches(6.0), Inches(2.95), Inches(5.55), Inches(0.9), fill=WHITE, line=LINE, rounded=True)
    text(slide, Inches(6.28), Inches(3.24), Inches(4.95), Inches(0.24), "不是“提示词写得漂亮”，而是“人先定边界，AI 再高效执行”。", size=13, bold=True)
    bullets(
        slide,
        Inches(6.12),
        Inches(4.28),
        Inches(5.1),
        Inches(1.2),
        ["方向未定时，AI 更适合暴露问题，不适合直接拍板。", "口径已定后，AI 最适合做同步、回写和一致性维护。"],
        size=12,
    )


def instruction(slide, total):
    base_bg(slide, 12, total, accent=CYAN, chapter="CHAPTER 6", title="协作指令模式")
    rect(slide, Inches(1.12), Inches(2.0), Inches(5.05), Inches(4.0), fill=SOFT_BLUE, line=LINE, rounded=True)
    rect(slide, Inches(6.78), Inches(2.0), Inches(5.04), Inches(4.0), fill=SOFT_RED, line=LINE, rounded=True)
    text(slide, Inches(1.36), Inches(2.28), Inches(2.0), Inches(0.2), "高价值说法", size=18, bold=True)
    text(slide, Inches(7.02), Inches(2.28), Inches(2.4), Inches(0.2), "容易带偏的说法", size=18, bold=True)
    bullets(
        slide,
        Inches(1.36),
        Inches(2.82),
        Inches(4.4),
        Inches(2.55),
        [
            "在我没有结束之前，你不用做任何分析，只接收和记录。",
            "结合代码和原始需求，对照 PRD 找问题。",
            "先不改，先告诉我你准备怎么改。",
            "手机端先不改，我们先改 PC 端，按照原型改。",
        ],
        size=13,
    )
    bullets(
        slide,
        Inches(7.02),
        Inches(2.82),
        Inches(4.4),
        Inches(2.55),
        [
            "你思考下有什么改进方案。",
            "你先把详情页原型改了，我看看效果。",
            "问题不在语气，而在边界和阶段不够清楚。",
        ],
        size=13,
    )
    rect(slide, Inches(1.12), Inches(5.95), Inches(10.7), Inches(0.4), fill=NAVY, rounded=True)
    text(slide, Inches(1.36), Inches(6.06), Inches(10.1), Inches(0.14), "好指令的核心不是更复杂，而是边界更清楚、阶段更明确。", size=12, bold=True, color=WHITE)


def conclusion(slide, total):
    base_bg(slide, 13, total, accent=GOLD, chapter="CHAPTER 7", title="最后想传递的结论")
    rect(slide, Inches(1.12), Inches(2.0), Inches(10.7), Inches(0.92), fill=NAVY, rounded=True)
    text(slide, Inches(1.42), Inches(2.28), Inches(10.1), Inches(0.24), "AI 的价值，不在于替产品经理做产品，而在于帮你更快地整理、比对、同步和暴露问题。", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(
        slide,
        Inches(1.36),
        Inches(3.55),
        Inches(9.9),
        Inches(2.15),
        [
            "AI 的提效是成立的，但这个提效有边界。",
            "真正决定产出质量的，仍然是产品经理有没有讲清边界、拍板规则、持续联动校正。",
            "AI 节省的是整理、核查、同步和回写时间，不替代业务判断。",
            "最值得沉淀的不是一组提示词，而是一套适用于存量功能迭代的协作方法。",
        ],
        size=14,
    )


def ending(slide, total):
    rect(slide, 0, 0, SW, SH, fill=NAVY)
    rect(slide, Inches(9.3), Inches(-0.5), Inches(4.0), Inches(2.4), fill=BLUE, oval=True)
    rect(slide, Inches(-0.6), Inches(5.7), Inches(3.1), Inches(1.9), fill=INDIGO, oval=True)
    rect(slide, Inches(0.95), Inches(1.0), Inches(11.4), Inches(5.45), fill=RGBColor(23, 45, 92), line=RGBColor(65, 91, 151), rounded=True)
    chip(slide, Inches(1.35), Inches(1.38), Inches(1.0), "ENDING", fill=WHITE, line=None, color=BLUE)
    text(
        slide,
        Inches(1.38),
        Inches(2.18),
        Inches(9.8),
        Inches(1.55),
        "在原有功能迭代项目里，\nAI 不是替代产品经理的设计者，\n而是放大产品经理能力的协作工具。",
        size=28,
        bold=True,
        color=WHITE,
    )
    bullets(
        slide,
        Inches(1.42),
        Inches(4.72),
        Inches(8.4),
        Inches(1.25),
        [
            "先定边界，再让 AI 展开",
            "先对照现实，再让 AI 输出",
            "先拍板口径，再让 AI 同步",
            "先控制范围，再让 AI 提效",
        ],
        size=14,
        color=CYAN,
    )
    text(slide, Inches(10.18), Inches(5.55), Inches(1.5), Inches(0.3), "谢谢", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    footer(slide, 14, total)


def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    pages = [
        cover,
        agenda,
        why,
        pain,
        ai_steps,
        ai_value,
        pitfall1,
        pitfall2,
        pitfall_common,
        success1,
        success2,
        instruction,
        conclusion,
        ending,
    ]
    total = len(pages)
    for fn in pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fn(slide, total)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
